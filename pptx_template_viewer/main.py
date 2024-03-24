import os
import sys

import pptx

def main():
    print(sys.argv)

    prs = pptx.Presentation(sys.argv[1])

    for index in range(len(prs.slide_layouts)):
        print(f'-----===== {index} =====-----')
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        for shape in slide.shapes:
            print(f'{shape.shape_type}')
            if shape.is_placeholder:
                print(f'[{shape.placeholder_format.idx}] {shape.placeholder_format.type}')
        print()


if __name__ == '__main__':
    main()

